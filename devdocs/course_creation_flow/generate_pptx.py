"""
StudyStudio 课程制作全流程 PPT 生成脚本
输出：course_creation_flow.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 配色方案 ──
BG_DARK   = RGBColor(0x1A, 0x1A, 0x2E)  # 深蓝黑背景
ACCENT    = RGBColor(0x00, 0xD2, 0xFF)  # 青色强调
ACCENT2   = RGBColor(0x7B, 0x2F, 0xBE)  # 紫色
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCC, 0xCC, 0xDD)
ORANGE    = RGBColor(0xFF, 0x8C, 0x00)
GREEN     = RGBColor(0x00, 0xE6, 0x76)
RED       = RGBColor(0xFF, 0x44, 0x44)
CARD_BG   = RGBColor(0x22, 0x22, 0x3A)
SUBTLE    = RGBColor(0x99, 0x99, 0xAA)

prs = Presentation()
prs.slide_width  = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

# ── 工具函数 ──

def add_bg(slide, color=BG_DARK):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, shape_type=MSO_SHAPE.RECTANGLE):
    """添加矩形"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_lines(slide, lines, left, top, width, height, font_size=16, color=WHITE,
              bold_first=False, spacing=Pt(26)):
    """添加多行文本"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        if bold_first and i == 0:
            p.font.bold = True
        p.space_after = spacing
    return txBox

def add_card(slide, left, top, width, height, number, title, desc, accent_color=ACCENT):
    """添加编号卡片"""
    card = add_rect(slide, left, top, width, height, CARD_BG)
    # 编号圆圈
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.15), top + Inches(0.15),
                                     Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent_color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    # 标题
    add_text_box(slide, left + Inches(0.8), top + Inches(0.15),
                 width - Inches(1.0), Inches(0.45), title, font_size=18, color=accent_color, bold=True)
    # 描述
    add_text_box(slide, left + Inches(0.2), top + Inches(0.75),
                 width - Inches(0.4), height - Inches(0.9), desc, font_size=12, color=LIGHT)
    return card

def slide_number(slide, num):
    """添加页码"""
    add_text_box(slide, Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.3),
                 str(num), font_size=10, color=SUBTLE, alignment=PP_ALIGN.RIGHT)

def accent_bar(slide, left, top, width=Inches(0.06), height=Inches(0.5), color=ACCENT):
    """添加竖条装饰"""
    return add_rect(slide, left, top, width, height, color)

# ═══════════════════════════════════════════
# Slide 1: 封面
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
add_bg(sl)
# 装饰条
add_rect(sl, Inches(0), Inches(3.0), Inches(13.333), Inches(0.04), ACCENT)
add_rect(sl, Inches(0), Inches(3.15), Inches(13.333), Inches(0.02), ACCENT2)
# 标题
add_text_box(sl, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             'StudyStudio 课程制作全流程', font_size=48, color=WHITE, bold=True)
add_text_box(sl, Inches(1), Inches(2.7), Inches(11), Inches(0.5),
             '从文件上传到学生学习的完整技术链路', font_size=22, color=ACCENT)
# 副标题
add_lines(sl, [
    '10 个自动化阶段 · 全异步事件驱动 · AI + 向量聚类合成',
    'Celery 任务编排 · RabbitMQ 消息队列 · MinIO 对象存储 · PostgreSQL 向量检索'
], Inches(1), Inches(3.6), Inches(11), Inches(1.5), font_size=16, color=LIGHT, spacing=Pt(20))
# 底部
add_text_box(sl, Inches(1), Inches(6.5), Inches(5), Inches(0.4),
             'StudyStudio DS  ·  v2.8.0  ·  2026-04', font_size=14, color=SUBTLE)

# ═══════════════════════════════════════════
# Slide 2: 总览
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_text_box(sl, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             '全流程总览：10 阶段流水线', font_size=32, color=WHITE, bold=True)
accent_bar(sl, Inches(0.5), Inches(0.5))

# 10 个阶段的概览 — 分两行 5+5
stages = [
    ('1', '文件上传', '格式校验\nSHA-256 去重\nMinIO 存储'),
    ('2', '事件调度', 'RabbitMQ 订阅\n→ ingest 队列'),
    ('3', '文档解析', 'PyMuPDF/pdfplumber\n分块 + 标准化'),
    ('4', '知识提取', '三步 LLM 流水线\n实体识别→分类→关系'),
    ('5', '自动审核', '两轮 LLM 评审\napprove/reject'),
    ('6', '向量化', 'Embedding 批量生成\npgvector 存储'),
    ('7', '蓝图生成', 'K-Means 聚类\nLLM 章节/内容生成'),
    ('8', '测验预生成', '每章自动出题\n写入题库'),
    ('9', '学生学习', '三种阅读模式\n测验+反思+笔记'),
    ('10', '教师精调', '自然语言指令\nLayer 2 精调联动'),
]
for i, (num, title, desc) in enumerate(stages):
    row = i // 5
    col = i % 5
    card_w = Inches(2.3)
    card_h = Inches(2.6)
    x = Inches(0.5 + col * 2.55)
    y = Inches(1.4 + row * 2.85)
    acc = ACCENT if row == 0 else ACCENT2
    add_card(sl, x, y, card_w, card_h, num, title, desc, acc)

# ═══════════════════════════════════════════
# Slide 3: 架构全景
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_text_box(sl, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             '系统架构全景', font_size=32, color=WHITE, bold=True)
accent_bar(sl, Inches(0.5), Inches(0.5))

arch_items = [
    ('前端 (Vue 3 + TS)', ACCENT, ['UploadView → 文件上传', 'TutorialView → 课程学习', 'API 客户端封装']),
    ('API 网关 (FastAPI)', ACCENT2, ['REST 端点 + WebSocket', 'Depends 依赖注入鉴权', 'EventBus 事件发布']),
    ('消息队列 (RabbitMQ)', ORANGE, ['file_uploaded → ingest', 'document_parsed → extraction', '蓝图/测验/讨论事件']),
    ('Celery Workers', GREEN, ['knowledge 队列：解析+提取+向量', 'knowledge.review 队列：自动审核', 'blueprint.synthesis 队列：生成']),
    ('存储层', RED, ['PostgreSQL + pgvector', 'MinIO 对象存储', 'Redis 结果/缓存']),
]

for i, (title, color, items) in enumerate(arch_items):
    x = Inches(0.4 + i * 2.55)
    y = Inches(1.5)
    # 标题
    add_text_box(sl, x, y, Inches(2.4), Inches(0.5), title, font_size=16, color=color, bold=True)
    y2 = y + Inches(0.6)
    # 卡片背景
    add_rect(sl, x, y2, Inches(2.4), Inches(4.5), CARD_BG)
    # 条目
    for j, item in enumerate(items):
        add_text_box(sl, x + Inches(0.15), y2 + Inches(0.3 + j * 0.55), Inches(2.1), Inches(0.5),
                     f'▸ {item}', font_size=12, color=LIGHT)
    # 箭头（除了最后一个）
    if i < len(arch_items) - 1:
        arrow = sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(2.45),
                                    Inches(5.8), Inches(0.3), Inches(0.2))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = SUBTLE
        arrow.line.fill.background()

# 底部状态机
add_text_box(sl, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
             '文档状态机：uploaded → parsed → extracting → extracted → embedding → reviewed → published',
             font_size=14, color=ORANGE, bold=True)
slide_number(sl, 3)

# ═══════════════════════════════════════════
# Slides 4-13: 各阶段详情
# ═══════════════════════════════════════════

stage_details = [
    {
        'num': 1,
        'title': '文件上传与存储',
        'subtitle': '前端 → API → MinIO，SHA-256 去重保证不重复处理',
        'icon': '☁',
        'left_points': [
            ('入口', 'UploadView.vue\n支持 PDF/DOCX/MD/TXT\n最大 100MB'),
            ('校验', 'Content-Type 白名单\n文件名安全过滤\n目录穿越防护'),
            ('去重', 'SHA-256 哈希计算\n同文件同空间 → 复用\n同文件不同空间 → 新记录'),
        ],
        'right_points': [
            ('存储', 'MinIO 对象存储\n键格式：files/{id}/{name}\nAsyncMinIOClient 异步封装'),
            ('事件', '发布 file_uploaded\n→ RabbitMQ\n→ knowledge.ingest.queue'),
            ('状态', 'documents 表\n初始状态：uploaded'),
        ],
        'code': 'POST /api/files/upload\n→ file_router.py\n→ sanitize + dedup\n→ MinIO put_object\n→ EventBus.publish("file_uploaded")',
    },
    {
        'num': 2,
        'title': '事件驱动任务调度',
        'subtitle': 'RabbitMQ 解耦，Celery 异步执行，长链路自动串联',
        'icon': '⚡',
        'left_points': [
            ('订阅注册', 'FastAPI startup 事件\nmain.py 注册 2 个订阅\nevent_bus.subscribe()'),
            ('路由规则', 'file_uploaded\n→ knowledge.ingest.queue\n→ run_ingest'),
            ('路由规则', 'document_parsed\n→ knowledge.extraction.queue\n→ run_extraction'),
        ],
        'right_points': [
            ('Celery 配置', 'Broker: RabbitMQ\nResult Backend: Redis\n5 个专用队列'),
            ('队列清单', 'knowledge\nknowledge.review\ntutorial\nblueprint.synthesis\nlow_priority'),
            ('超时策略', 'ingest: 180s soft / 240s hard\nextraction: 无限制\nreview: 600s soft / 720s hard'),
        ],
        'code': '@event_bus.subscribe("file_uploaded")\nasync def on_file_uploaded(event):\n    run_ingest.delay(event.data)\n    logger.info("Dispatched ingest task")',
    },
    {
        'num': 3,
        'title': '文档解析与文本分块',
        'subtitle': '多格式解析引擎 + RecursiveCharacterTextSplitter',
        'icon': '📄',
        'left_points': [
            ('PDF 解析', 'PyMuPDF (fitz)\n按物理页提取文本\nTOC 大纲提取'),
            ('DOCX 解析', 'python-docx\n段落级提取\n保留结构层次'),
            ('MD/TXT', 'charset 检测 → 解码\n不依赖文件扩展名'),
        ],
        'right_points': [
            ('文本分块', 'LangChain Recursive\nCharacterTextSplitter\n每页独立分块'),
            ('元数据', 'index_no / title_path\npage_no / token_count\n全文搜索友好'),
            ('截断保护', 'MAX_CHUNK_COUNT = 500\n超出标记 is_truncated\n写入 document_chunks 表'),
        ],
        'code': 'DocumentIngestService.ingest()\n→ minio.download(tmp_path)\n→ _extract_pages() per format\n→ _split_text() per page\n→ batch INSERT chunks (50/batch)\n→ status = "parsed"\n→ emit "document_parsed"',
    },
    {
        'num': 4,
        'title': '知识提取 — 三步 LLM 流水线',
        'subtitle': 'Step 1 实体识别 → Step 2 分类定义 → Step 3 关系抽取',
        'icon': '🧠',
        'left_points': [
            ('Step 1: 识别', 'ENTITY_RECOGNITION_PROMPT\n从文本提取教学实体\n输出实体名称列表'),
            ('Step 2: 分类', 'ENTITY_CLASSIFICATION_PROMPT\n5 种类别:\nconcept/element/flow/case/defense\n附加 short_definition'),
            ('Step 3: 关系', 'RELATION_EXTRACTION_PROMPT\n4 种关系:\nprerequisite_of/related_to\npart_of/example_of'),
        ],
        'right_points': [
            ('跨块去重', '按 entity_name 合并\nall_entities 跨块累积'),
            ('跨文档去重', '同 canonical_name\n+ 同 domain_tag → 跳过'),
            ('批量写入', 'knowledge_entities\nreview_status = "pending"\nknowledge_relations'),
            ('原子锁', 'UPDATE documents SET\nstatus = "extracting"\nWHERE NOT IN(extracting,extracted)'),
        ],
        'code': 'run_extraction(chunk_data)\n→ _step_entity_recognition()\n→ _step_entity_classification()\n→ _step_relation_extraction()\n→ cross-chunk dedup\n→ INSERT entities + relations\n→ status = "extracted"\n→ dispatch auto_review_entities',
    },
    {
        'num': 5,
        'title': 'AI 自动审核 — 两轮质量把关',
        'subtitle': '保证无实体永久 pending，自动决定 approve/reject',
        'icon': '✅',
        'left_points': [
            ('Round 1', 'LLM 逐实体评审\n输出: approve/reject/uncertain\n含 confidence 评分'),
            ('裁决规则', 'approve + confidence ≥ 0.75\n→ "approved"\nreject → "rejected"\n其余 → Round 2'),
            ('Round 2', '更严格二次审核\napprove + confidence ≥ 0.60\n→ "approved"\n其余 → "rejected"'),
        ],
        'right_points': [
            ('并发安全', 'pg_try_advisory_lock()\n空间级锁\n5 entity/batch'),
            ('守护任务', 'resume_pending_review\n每 5 分钟 celery_beat\n抢救卡住的 review'),
            ('完成动作', '→ 文档状态 → "embedding"\n→ 触发向量化任务'),
        ],
        'code': 'auto_review_entities(space_id)\n→ lock via advisory lock\n→ for each entity batch:\n    ROUND1_PROMPT → LLM\n    分流: approved/rejected/round2\n    ROUND2_PROMPT → LLM\n→ _finalize_documents()\n→ status = "embedding"',
    },
    {
        'num': 6,
        'title': '向量化 — Embedding 生成',
        'subtitle': '批量向量化 + pgvector 存储 + 相似度检索',
        'icon': '🔢',
        'left_points': [
            ('嵌入文本', 'canonical_name + "—"\n+ short_definition\n截断 512 字符'),
            ('批量生成', 'backfill_entity_embeddings\nbatch_size = 32\nLLMGateway.embed()'),
            ('存储格式', 'PostgreSQL vector 类型\npgvector 扩展\n索引加速'),
        ],
        'right_points': [
            ('Chunk 向量', 'embed_document_chunks\n并行执行\n用于原文相似度搜索'),
            ('完成检测', '_trigger_blueprint_if_ready()\n检查无 pending embedding\n→ status → "reviewed"'),
            ('触发合成', '→ dispatch\nsynthesize_blueprint\n→ blueprint.synthesis.queue'),
        ],
        'code': 'backfill_entity_embeddings(space_id)\n→ SELECT WHERE embedding IS NULL\n→ _build_embed_text() per entity\n→ batch embed via LLMGateway\n→ UPDATE embedding = CAST(:emb AS vector)\n→ _trigger_blueprint_if_ready()',
    },
    {
        'num': 7,
        'title': '蓝图/课程生成 — V2 合成引擎',
        'subtitle': 'K-Means 聚类 + LLM 多阶段生成 (全新 / 增量合并)',
        'icon': '🏗',
        'left_points': [
            ('过滤', '_is_non_teaching_entity()\n过滤 CVE/版本号/厂商名\n保留教学价值实体'),
            ('聚类', 'K-Means on embedding\nK = max(4, round(N/7))\n上限 30\n+ _rebalance_clusters()'),
            ('LLM 生成 4 步', '① _name_cluster() 章节命名\n② _plan_stages_with_llm() 阶段规划\n③ COURSE_TITLE_PROMPT 课程标题\n④ _gen_content() 章节正文'),
        ],
        'right_points': [
            ('模式 A: 全新', '首次生成完整课程\n3 阶段 (基础+实践+评估)\n每章 600-900 字\n结构化: scene_hook +\nskim_summary + code_example\n+ misconception_block'),
            ('模式 B: 合并', '追加文档到已有课程\n余弦相似度差异分析\n>0.92 跳过 <0.75 新章\n进度自动重映射'),
            ('状态', 'author_status = "published"\ndocuments → "published"'),
        ],
        'code': 'synthesize_blueprint()\n→ _acquire_or_wait_blueprint_lock()\n→ if first: _synthesize_blueprint_v2_async()\n   elif merge: _synthesize_blueprint_merge_async()\n→ write blueprint + stages + chapters\n→ dispatch pregen_chapter_quizzes',
    },
    {
        'num': 8,
        'title': 'Layer 1 — 生成前约束 (Teacher Instruction)',
        'subtitle': '教师在上传时指定全局教学约束，影响整个生成链路',
        'icon': '🎯',
        'left_points': [
            ('输入时机', '文件上传时\n或课程列表页\n"自定义提示词" 输入框'),
            ('注入位置', '写入 skill_blueprints\n.teacher_instruction\n在所有 LLM Prompt 中拼接'),
            ('影响范围', '章节标题生成\n内容风格控制\n难度 / 受众适配'),
        ],
        'right_points': [
            ('示例指令', '"增加实操案例，弱化理论推导"\n"适配中职三年级基础水平"\n"加入航空维修安全规范"\n"用 Java 代替 Python 示例"'),
            ('存储字段', 'teacher_instruction TEXT\nversion 字段关联\n每个蓝图版本独立记录'),
            ('联动', 'Layer 2 精调可覆盖/\n补充 Layer 1 约束'),
        ],
        'code': 'if teacher_instruction:\n    prompt += f"""\n    教学约束提示：\n    {teacher_instruction}\n    """\n蓝图生成时注入 CLUSTER_CHAPTER_PROMPT\n和 CHAPTER_CONTENT_PROMPT',
    },
    {
        'num': 9,
        'title': '学生学习与交互',
        'subtitle': '三种阅读模式 + 测验反思 + 社交学习 + 证书',
        'icon': '📖',
        'left_points': [
            ('阅读模式', '① 速览: skim_summary 要点\n② 正常: 完整正文\n③ 深度: + prereq_adaptive'),
            ('内容结构', 'scene_hook 场景导入\nfull_content 正文\nCHECKPOINT 交互断点\ncode_example 代码示例\nmisconception_block 误区'),
            ('测验系统', '/chapter-quiz/{id}\n随机抽题\n≥ 60% 自动标记已读'),
        ],
        'right_points': [
            ('AI 反思', '学生写理解总结\nLLM 评分反馈\n一对一批改'),
            ('社交笔记', '每章独立笔记区\n点赞互动\n同伴学习'),
            ('其他', '源文档追溯\n关联推荐\n证书下载 (全部已读)\n蓝图更新订阅通知'),
        ],
        'code': 'TutorialView.vue\n→ loadTutorial() 加载课程\n→ chapterContent computed 解析\n→ contentSegments 切分 CHECKPOINT\n→ quiz / reflection / notes\n→ markChapterRead()',
    },
    {
        'num': 10,
        'title': 'Layer 2+3 — 教师精调与联动再生',
        'subtitle': '自然语言指令 → 章节重写 → 测验/讨论联动更新',
        'icon': '✨',
        'left_points': [
            ('精调入口', 'TutorialView 章节标题旁\n"✨ 精调本章" 按钮\nElMessageBox textarea\n自然语言描述要求'),
            ('后端处理', '读取当前章节元数据\n构建 CHAPTER_REFINEMENT_PROMPT\n含当前内容摘要\n+ teacher_instruction\n+ 全局约束'),
            ('LLM 调用', 'route: tutorial_content\ntimeout: 150s\n生成新章节完整内容'),
        ],
        'right_points': [
            ('Layer 3 联动', '① 测验缓存失效\n② regenerate_chapter_quiz\n    重新生成题目\n③ generate_discussion_seeds\n    可选讨论种子生成'),
            ('版本管理', 'refinement_version + 1\n保留精调历史'),
            ('鉴权', 'require_space_owner()\n仅课程所有者可用\n非系统管理员'),
        ],
        'code': 'POST /admin/courses/chapters/{id}/refine\n→ get_current_user\n→ require_space_owner(space_id)\n→ build CHAPTER_REFINEMENT_PROMPT\n→ LLM generate\n→ update_chapter_content()\n→ invalidate_quiz + regenerate',
    },
]

for detail in stage_details:
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(sl)

    # 顶部：阶段编号 + 标题
    add_text_box(sl, Inches(0.8), Inches(0.2), Inches(0.8), Inches(0.7),
                 f'{detail["icon"]}', font_size=36, color=WHITE)
    add_text_box(sl, Inches(1.4), Inches(0.2), Inches(10), Inches(0.8),
                 f'阶段 {detail["num"]}：{detail["title"]}', font_size=30, color=WHITE, bold=True)
    add_text_box(sl, Inches(1.4), Inches(0.85), Inches(10), Inches(0.4),
                 detail['subtitle'], font_size=14, color=ACCENT)
    accent_bar(sl, Inches(0.5), Inches(0.3))

    # 分隔线
    add_rect(sl, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.02), ACCENT)

    # 左侧要点
    for j, (label, text) in enumerate(detail['left_points']):
        y = Inches(1.55 + j * 1.65)
        add_text_box(sl, Inches(0.6), y, Inches(1.0), Inches(0.3),
                     label, font_size=13, color=ORANGE, bold=True)
        add_rect(sl, Inches(1.65), y + Inches(0.05), Inches(4.8), Inches(1.45), CARD_BG)
        add_text_box(sl, Inches(1.8), y + Inches(0.1), Inches(4.5), Inches(1.3),
                     text, font_size=12, color=LIGHT)

    # 右侧要点
    for j, (label, text) in enumerate(detail['right_points']):
        y = Inches(1.55 + j * 1.65)
        add_text_box(sl, Inches(6.7), y, Inches(1.0), Inches(0.3),
                     label, font_size=13, color=GREEN, bold=True)
        add_rect(sl, Inches(7.75), y + Inches(0.05), Inches(5.1), Inches(1.45), CARD_BG)
        add_text_box(sl, Inches(7.9), y + Inches(0.1), Inches(4.8), Inches(1.3),
                     text, font_size=12, color=LIGHT)

    # 底部代码
    add_rect(sl, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.95), RGBColor(0x0D, 0x0D, 0x15))
    add_text_box(sl, Inches(0.7), Inches(6.45), Inches(11.9), Inches(0.85),
                 detail['code'], font_size=11, color=GREEN)

    slide_number(sl, 3 + detail['num'])

# ═══════════════════════════════════════════
# Slide 14: 状态机全景
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_text_box(sl, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             '文档状态机 & 事件链路', font_size=32, color=WHITE, bold=True)
accent_bar(sl, Inches(0.5), Inches(0.5))

states = [
    ('uploaded', 'orange', 'run_ingest'),
    ('parsed', 'blue', 'run_extraction'),
    ('extracting', 'yellow', '原子锁'),
    ('extracted', 'purple', 'auto_review'),
    ('embedding', 'cyan', 'backfill_embed'),
    ('reviewed', 'green', 'synthesize'),
    ('published', 'bright', '完成'),
]

for i, (state, color_name, trigger) in enumerate(states):
    x = Inches(0.6 + i * 1.78)
    # 状态框
    box = add_rect(sl, x, Inches(2.0), Inches(1.6), Inches(1.2), CARD_BG)
    add_text_box(sl, x + Inches(0.1), Inches(2.1), Inches(1.4), Inches(0.5),
                 state, font_size=18, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.1), Inches(2.6), Inches(1.4), Inches(0.4),
                 trigger, font_size=11, color=SUBTLE, alignment=PP_ALIGN.CENTER)
    # 箭头
    if i < len(states) - 1:
        arrow = sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(1.6),
                                    Inches(2.45), Inches(0.18), Inches(0.15))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT
        arrow.line.fill.background()

# 事件流
add_text_box(sl, Inches(0.8), Inches(3.8), Inches(11), Inches(0.4),
             'RabbitMQ 事件链路', font_size=20, color=WHITE, bold=True)

events_flow = [
    ('file_uploaded', 'knowledge.ingest.queue', ACCENT),
    ('document_parsed', 'knowledge.extraction.queue', ACCENT),
    ('entities_reviewed', '(内部调用)', ACCENT2),
    ('embeddings_complete', '(内部调用)', ACCENT2),
    ('blueprint_published', 'blueprint.notify', ORANGE),
]

for i, (event, queue, color) in enumerate(events_flow):
    y = Inches(4.4 + i * 0.55)
    add_rect(sl, Inches(0.8), y, Inches(4.5), Inches(0.4), CARD_BG)
    add_text_box(sl, Inches(0.9), y + Inches(0.05), Inches(4.3), Inches(0.3),
                 f'事件: {event}', font_size=13, color=color)
    arrow2 = sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.3), y + Inches(0.07),
                                  Inches(0.2), Inches(0.2))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = SUBTLE
    arrow2.line.fill.background()
    add_text_box(sl, Inches(5.6), y + Inches(0.05), Inches(6), Inches(0.3),
                 f'→ {queue}', font_size=13, color=LIGHT)

slide_number(sl, 14)

# ═══════════════════════════════════════════
# Slide 15: Layer 概览
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_text_box(sl, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             '教师引导式迭代 — 三层方案', font_size=32, color=WHITE, bold=True)
accent_bar(sl, Inches(0.5), Inches(0.5))

layers = [
    ('Layer 1', '生成前约束', ACCENT, [
        '上传时输入全局教学约束',
        '注入所有 LLM Prompt',
        '控制风格/难度/受众',
    ]),
    ('Layer 2', '对话式精调', ACCENT2, [
        '自然语言指令重写章节',
        '✨ 精调本章 按钮',
        '保留 refinement_version 历史',
    ]),
    ('Layer 3', '附属内容联动', ORANGE, [
        '测验自动重新生成',
        '讨论种子可选更新',
        '缓存自动失效',
    ]),
]

for i, (layer, title, color, items) in enumerate(layers):
    x = Inches(0.8 + i * 4.1)
    # 层级标签
    tag = add_rect(sl, x, Inches(1.6), Inches(3.5), Inches(0.6), color)
    add_text_box(sl, x + Inches(0.1), Inches(1.65), Inches(3.3), Inches(0.5),
                 f'{layer}: {title}', font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # 内容卡片
    add_rect(sl, x, Inches(2.3), Inches(3.5), Inches(2.5), CARD_BG)
    for j, item in enumerate(items):
        add_text_box(sl, x + Inches(0.2), Inches(2.5 + j * 0.7), Inches(3.1), Inches(0.6),
                     f'✓ {item}', font_size=14, color=LIGHT)
    # 向下箭头
    if i < len(layers) - 1:
        arrow = sl.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(0.8 + i * 4.1 + 1.5),
                                     Inches(5.0), Inches(0.2), Inches(0.2))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = SUBTLE
        arrow.line.fill.background()

# 底部说明
add_text_box(sl, Inches(0.8), Inches(5.5), Inches(11), Inches(1.5),
             '鉴权模型 (v2.8.0 修复)：课程管理使用 require_space_owner() — 仅课程所有者（教师）可用，\n'
             '系统管理员 (admin 角色) 管理的是 AI 配置、用户管理、权限控制等系统级功能。\n'
             '管理端点：POST /admin/courses/chapters/{id}/refine（精调）、regen、reorder、delete、transfer',
             font_size=13, color=SUBTLE)

slide_number(sl, 15)

# ═══════════════════════════════════════════
# Slide 16: 技术栈汇总
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_text_box(sl, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             '技术栈一览', font_size=32, color=WHITE, bold=True)
accent_bar(sl, Inches(0.5), Inches(0.5))

tech_stacks = [
    ('前端', 'Vue 3 + TypeScript\nVite\nElement Plus\nPinia'),
    ('API', 'FastAPI\nPydantic v2\nUvicorn\nWebSocket'),
    ('任务', 'Celery\nRabbitMQ\nRedis\nEventBus'),
    ('AI', 'LangChain\nLLMGateway\nPyMuPDF\npython-docx'),
    ('数据', 'PostgreSQL\npgvector\nSQLAlchemy 2.0\nAlembic'),
    ('存储', 'MinIO (S3)\nboto3\nDocker\nDocker Compose'),
]

for i, (cat, items) in enumerate(tech_stacks):
    x = Inches(0.4 + i * 2.15)
    add_text_box(sl, x, Inches(1.5), Inches(2.0), Inches(0.5),
                 cat, font_size=20, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_rect(sl, x, Inches(2.1), Inches(2.0), Inches(1.8), CARD_BG)
    add_text_box(sl, x + Inches(0.15), Inches(2.2), Inches(1.7), Inches(1.5),
                 items, font_size=14, color=LIGHT, alignment=PP_ALIGN.CENTER)

# 指标
add_text_box(sl, Inches(0.8), Inches(4.3), Inches(11), Inches(0.5),
             '项目规模指标', font_size=22, color=WHITE, bold=True)

metrics = [
    ('27', '数据库迁移'),
    ('144', '测试用例'),
    ('9', 'Docker 服务'),
    ('10', '流水线阶段'),
    ('5', 'Celery 队列'),
    ('3', '教师迭代层'),
]
for i, (num, label) in enumerate(metrics):
    x = Inches(0.7 + i * 2.1)
    add_text_box(sl, x, Inches(5.0), Inches(1.8), Inches(0.8),
                 num, font_size=42, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x, Inches(5.8), Inches(1.8), Inches(0.4),
                 label, font_size=14, color=SUBTLE, alignment=PP_ALIGN.CENTER)

slide_number(sl, 16)

# ═══════════════════════════════════════════
# Slide 17: 结束页
# ═══════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
add_rect(sl, Inches(0), Inches(3.0), Inches(13.333), Inches(0.04), ACCENT)
add_rect(sl, Inches(0), Inches(3.15), Inches(13.333), Inches(0.02), ACCENT2)
add_text_box(sl, Inches(1), Inches(1.8), Inches(11), Inches(1.0),
             '谢谢', font_size=56, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(sl, Inches(1), Inches(3.6), Inches(11), Inches(0.8),
             'StudyStudio DS  ·  v2.8.0  ·  全异步 AI 驱动课程生成系统',
             font_size=18, color=LIGHT, alignment=PP_ALIGN.CENTER)
add_text_box(sl, Inches(1), Inches(4.5), Inches(11), Inches(0.8),
             'devdocs/course_creation_flow/ — PPT · Markdown · AI 提示词 · 视频脚本',
             font_size=14, color=SUBTLE, alignment=PP_ALIGN.CENTER)

# ── 保存 ──
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, 'StudyStudio_课程制作全流程.pptx')
prs.save(output_path)
print(f'PPT generated: {output_path}')
print(f'Total slides: {len(prs.slides)}')
